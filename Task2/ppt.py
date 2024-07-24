from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Function to add image to PowerPoint
def add_image_to_slide(slide, image_path, title, left, top, width):
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_path, left, top, width)

# Create a PowerPoint presentation
prs = Presentation()

# Add a title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Model Evaluation Results"
subtitle.text = "Random Forest Classifier on Customer Booking Data"

# Add ROC Curve slide
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
add_image_to_slide(slide, 'roc_curve.png', 'ROC Curve', Inches(1), Inches(1), Inches(8.5))

# Add Feature Importance slide
slide = prs.slides.add_slide(slide_layout)
add_image_to_slide(slide, 'feature_importances.png', 'Feature Importances', Inches(1), Inches(1), Inches(8.5))

# Save the presentation
prs.save('model_evaluation_results.pptx')

print("Presentation generated successfully.")