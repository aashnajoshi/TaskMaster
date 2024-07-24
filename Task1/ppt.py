from pptx import Presentation
from pptx.util import Inches

# Load the template
prs = Presentation('data/Presentation Template - Task 1.pptx')

# Access the slide
slide = prs.slides[0]

# Add title and subtitle
title = slide.shapes.title
title.text = "British Airways Customer Feedback Analysis"

# Add sentiment analysis chart
slide.shapes.add_picture('data/Figure 1.png', Inches(1), Inches(2), width=Inches(5))

# Add word cloud
slide.shapes.add_picture('data/Figure 2.png', Inches(6), Inches(2), width=Inches(5))

# Add key findings
text_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(10), Inches(1.5))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "Key Findings:"
p = text_frame.add_paragraph()
p.text = "1. Positive reviews: 60%"
p = text_frame.add_paragraph()
p.text = "2. Negative reviews: 25%"
p = text_frame.add_paragraph()
p.text = "3. Neutral reviews: 15%"
p = text_frame.add_paragraph()
p.text = "4. Common topics: service quality, punctuality, comfort"

# Save the presentation
prs.save('BA_Customer_Feedback_Analysis.pptx')
